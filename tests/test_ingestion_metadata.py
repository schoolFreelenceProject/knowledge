import hashlib
from io import BytesIO

import pytest
from docx import Document as DocxDocument
from openpyxl import Workbook
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from sqlalchemy import create_engine, select
from sqlalchemy.orm import sessionmaker

from app.db.models import (
    Base,
    DocumentChunkRecord,
    DocumentPermissionRecord,
    DocumentRecord,
    DocumentStatus,
    UserRecord,
)
from app.schemas.documents import EmbeddedChunk
from app.services.ingestion_service import FolderUploadItem, IngestionService
from app.services.metadata_service import (
    DocumentMetadataService,
    MetadataPersistenceError,
)
from app.services.permission_service import PermissionPersistenceError, PermissionService
from app.services.text_chunker import ChunkingConfig
from app.services.vector_store import StoredVectorBatch, build_point_id


class FakeEmbeddingService:
    def embed_chunks(self, chunks):
        return [
            EmbeddedChunk(
                vector=[float(index), 0.0, 1.0],
                text=chunk.text,
                metadata=chunk.metadata,
            )
            for index, chunk in enumerate(chunks, start=1)
        ]


class FakeVectorStore:
    def __init__(self) -> None:
        self.deleted_point_ids: list[str] = []

    def store_embeddings(self, embedded_chunks):
        embedded_chunk_list = list(embedded_chunks)
        return StoredVectorBatch(
            collection_name="company_documents",
            stored_count=len(embedded_chunk_list),
            vector_size=3,
            point_ids=[
                f"point-{chunk.metadata.chunk_index}"
                for chunk in embedded_chunk_list
            ],
        )

    def delete_points(self, point_ids):
        self.deleted_point_ids.extend(point_ids)


class TrackingVectorStore:
    collection_name = "company_documents"

    def __init__(self) -> None:
        self.deleted_point_ids: list[str] = []
        self.stored_point_ids: list[str] = []
        self.active_point_ids: set[str] = set()

    def store_embeddings(self, embedded_chunks):
        embedded_chunk_list = list(embedded_chunks)
        point_ids = [build_point_id(chunk) for chunk in embedded_chunk_list]
        self.stored_point_ids.extend(point_ids)
        self.active_point_ids.update(point_ids)
        return StoredVectorBatch(
            collection_name=self.collection_name,
            stored_count=len(embedded_chunk_list),
            vector_size=3,
            point_ids=point_ids,
        )

    def delete_points(self, point_ids):
        point_id_list = list(point_ids)
        self.deleted_point_ids.extend(point_id_list)
        self.active_point_ids.difference_update(point_id_list)


class FailingMetadataService:
    def save_document_metadata(self, **_kwargs):
        raise MetadataPersistenceError("simulated PostgreSQL transaction failure")


class FailingForPathMetadataService(DocumentMetadataService):
    def __init__(self, *args, fail_source_path: str, **kwargs) -> None:
        super().__init__(*args, **kwargs)
        self.fail_source_path = fail_source_path

    def save_document_metadata(self, extracted_documents, **kwargs):
        if extracted_documents[0].metadata.source_path == self.fail_source_path:
            raise MetadataPersistenceError("simulated PostgreSQL transaction failure")

        return super().save_document_metadata(
            extracted_documents=extracted_documents,
            **kwargs,
        )


class FakePermissionService:
    def __init__(self, should_fail: bool = False) -> None:
        self.should_fail = should_fail
        self.grants: list[tuple[int, int]] = []

    def grant_document_access(self, document_id: int, user_id: int):
        if self.should_fail:
            raise PermissionPersistenceError("simulated permission failure")

        self.grants.append((document_id, user_id))


def _build_sqlite_metadata_service():
    engine = create_engine("sqlite:///:memory:")
    session_factory = sessionmaker(bind=engine, expire_on_commit=False)
    metadata_service = DocumentMetadataService(
        session_factory=session_factory,
        init_database=lambda: Base.metadata.create_all(bind=engine),
    )
    return metadata_service, session_factory


def _build_folder_services(
    tmp_path,
    max_upload_bytes: int | None = None,
    metadata_service: DocumentMetadataService | None = None,
    vector_store: TrackingVectorStore | None = None,
):
    engine = create_engine("sqlite:///:memory:")
    Base.metadata.create_all(bind=engine)
    session_factory = sessionmaker(bind=engine, expire_on_commit=False)
    resolved_metadata_service = metadata_service or DocumentMetadataService(
        session_factory=session_factory,
        init_database=lambda: Base.metadata.create_all(bind=engine),
    )
    permission_service = PermissionService(
        session_factory=session_factory,
        init_database=lambda: Base.metadata.create_all(bind=engine),
    )
    with session_factory() as session:
        session.add(
            UserRecord(
                id=7,
                email="uploader@example.com",
                password_hash="$argon2id$hash",
            )
        )
        session.add(
            UserRecord(
                id=8,
                email="second@example.com",
                password_hash="$argon2id$hash",
            )
        )
        session.commit()

    resolved_vector_store = vector_store or TrackingVectorStore()
    ingestion_service = IngestionService(
        documents_dir=tmp_path,
        chunk_config=ChunkingConfig(chunk_size=200, chunk_overlap=20),
        embedding_service=FakeEmbeddingService(),
        vector_store=resolved_vector_store,
        metadata_service=resolved_metadata_service,
        permission_service=permission_service,
        max_upload_bytes=max_upload_bytes,
    )
    return (
        ingestion_service,
        resolved_metadata_service,
        permission_service,
        resolved_vector_store,
        session_factory,
    )


def test_ingestion_metadata_persistence(tmp_path) -> None:
    metadata_service, session_factory = _build_sqlite_metadata_service()
    content = b"# Security Policy\n\nEmployees must report incidents quickly."

    ingestion_service = IngestionService(
        documents_dir=tmp_path,
        chunk_config=ChunkingConfig(chunk_size=200, chunk_overlap=20),
        embedding_service=FakeEmbeddingService(),
        vector_store=FakeVectorStore(),
        metadata_service=metadata_service,
    )

    response = ingestion_service.ingest_uploaded_document(
        filename="security.md",
        content=content,
    )

    with session_factory() as session:
        document = session.scalars(select(DocumentRecord)).one()
        chunks = session.scalars(select(DocumentChunkRecord)).all()

    assert response.document_id == document.id
    assert response.status == DocumentStatus.INDEXED.value
    assert response.file_hash == hashlib.sha256(content).hexdigest()
    assert document.storage_path == "security.md"
    assert document.file_hash == response.file_hash
    assert document.status == DocumentStatus.INDEXED.value
    assert len(chunks) == response.saved_chunks == response.chunks
    assert chunks[0].qdrant_point_id == "point-1"


def test_metadata_failure_cleans_up_qdrant_points_and_uploaded_file(tmp_path) -> None:
    vector_store = FakeVectorStore()
    ingestion_service = IngestionService(
        documents_dir=tmp_path,
        chunk_config=ChunkingConfig(chunk_size=200, chunk_overlap=20),
        embedding_service=FakeEmbeddingService(),
        vector_store=vector_store,
        metadata_service=FailingMetadataService(),
    )

    with pytest.raises(MetadataPersistenceError):
        ingestion_service.ingest_uploaded_document(
            filename="security.md",
            content=b"# Security Policy\n\nEmployees must report incidents quickly.",
        )

    assert vector_store.deleted_point_ids == ["point-1"]
    assert not (tmp_path / "security.md").exists()


def test_ingestion_auto_grants_uploader_document_access(tmp_path) -> None:
    metadata_service, session_factory = _build_sqlite_metadata_service()
    permission_service = FakePermissionService()
    ingestion_service = IngestionService(
        documents_dir=tmp_path,
        chunk_config=ChunkingConfig(chunk_size=200, chunk_overlap=20),
        embedding_service=FakeEmbeddingService(),
        vector_store=FakeVectorStore(),
        metadata_service=metadata_service,
        permission_service=permission_service,
    )

    response = ingestion_service.ingest_uploaded_document(
        filename="security.md",
        content=b"# Security Policy\n\nEmployees must report incidents quickly.",
        uploader_user_id=7,
    )

    assert permission_service.grants == [(response.document_id, 7)]
    with session_factory() as session:
        assert session.get(DocumentRecord, response.document_id) is not None


def test_single_japanese_docx_upload_indexes_and_preserves_acl_and_filename(
    tmp_path,
) -> None:
    (
        ingestion_service,
        _metadata_service,
        permission_service,
        _vector_store,
        session_factory,
    ) = _build_folder_services(tmp_path)
    content = _docx_bytes("休暇制度", "日本語の有給休暇ポリシーです。")

    response = ingestion_service.ingest_uploaded_document(
        filename="就業規則.docx",
        content=content,
        uploader_user_id=7,
    )

    with session_factory() as session:
        document = session.get(DocumentRecord, response.document_id)
        chunks = session.scalars(select(DocumentChunkRecord)).all()

    assert response.file_type == "docx"
    assert response.filename == "就業規則.docx"
    assert document is not None
    assert document.storage_path == "就業規則.docx"
    assert permission_service.list_accessible_document_ids(7) == [response.document_id]
    assert chunks[0].section_heading == "休暇制度"


def test_ingestion_refreshes_retrieval_index_after_successful_upload(tmp_path) -> None:
    metadata_service, _session_factory = _build_sqlite_metadata_service()
    refresh_calls: list[str] = []
    ingestion_service = IngestionService(
        documents_dir=tmp_path,
        chunk_config=ChunkingConfig(chunk_size=200, chunk_overlap=20),
        embedding_service=FakeEmbeddingService(),
        vector_store=FakeVectorStore(),
        metadata_service=metadata_service,
        retrieval_index_refresh=lambda: refresh_calls.append("refresh"),
    )

    ingestion_service.ingest_uploaded_document(
        filename="security.md",
        content=b"# Security Policy\n\nEmployees must report incidents quickly.",
    )

    assert refresh_calls == ["refresh"]


def test_permission_failure_cleans_up_ingested_state(tmp_path) -> None:
    metadata_service, session_factory = _build_sqlite_metadata_service()
    vector_store = FakeVectorStore()
    ingestion_service = IngestionService(
        documents_dir=tmp_path,
        chunk_config=ChunkingConfig(chunk_size=200, chunk_overlap=20),
        embedding_service=FakeEmbeddingService(),
        vector_store=vector_store,
        metadata_service=metadata_service,
        permission_service=FakePermissionService(should_fail=True),
    )

    with pytest.raises(MetadataPersistenceError):
        ingestion_service.ingest_uploaded_document(
            filename="security.md",
            content=b"# Security Policy\n\nEmployees must report incidents quickly.",
            uploader_user_id=7,
        )

    with session_factory() as session:
        assert session.query(DocumentRecord).count() == 0
        assert session.query(DocumentChunkRecord).count() == 0

    assert vector_store.deleted_point_ids == ["point-1"]
    assert not (tmp_path / "security.md").exists()


def test_folder_ingestion_indexes_nested_documents_preserves_paths_and_acl(
    tmp_path,
) -> None:
    (
        ingestion_service,
        _metadata_service,
        _permission_service,
        vector_store,
        session_factory,
    ) = _build_folder_services(tmp_path)

    response = ingestion_service.ingest_folder_documents(
        folder_name="CompanyDocs",
        files=[
            FolderUploadItem(
                relative_path="HR/leave.md",
                content=b"# Leave Policy\n\nEmployees can request paid leave.",
            ),
            FolderUploadItem(
                relative_path="IT/security.md",
                content=b"# Security Policy\n\nUse hardware keys for admin access.",
            ),
        ],
        uploader_user_id=7,
    )

    with session_factory() as session:
        documents = session.scalars(
            select(DocumentRecord).order_by(DocumentRecord.storage_path)
        ).all()
        chunks = session.scalars(select(DocumentChunkRecord)).all()
        permissions = session.scalars(
            select(DocumentPermissionRecord).order_by(
                DocumentPermissionRecord.document_id
            )
        ).all()

    assert response.folder_name == "CompanyDocs"
    assert response.files_discovered == 2
    assert response.indexed == 2
    assert response.skipped == 0
    assert response.failed == 0
    assert [document.filename for document in documents] == [
        "HR/leave.md",
        "IT/security.md",
    ]
    assert [document.storage_path for document in documents] == [
        "HR/leave.md",
        "IT/security.md",
    ]
    assert [(permission.document_id, permission.user_id) for permission in permissions] == [
        (documents[0].id, 7),
        (documents[1].id, 7),
    ]
    assert len(
        {(chunk.document_id, chunk.chunk_index) for chunk in chunks}
    ) == len(chunks)
    assert vector_store.active_point_ids == {
        chunk.qdrant_point_id for chunk in chunks
    }
    assert (tmp_path / "HR" / "leave.md").exists()
    assert (tmp_path / "IT" / "security.md").exists()


def test_folder_ingestion_indexes_mixed_formats_and_continues_after_bad_office_file(
    tmp_path,
) -> None:
    (
        ingestion_service,
        _metadata_service,
        _permission_service,
        vector_store,
        session_factory,
    ) = _build_folder_services(tmp_path)

    response = ingestion_service.ingest_folder_documents(
        folder_name="MixedDocs",
        files=[
            FolderUploadItem(
                relative_path="Policies/leave.md",
                content=b"# Leave\n\nEmployees can request paid leave.",
            ),
            FolderUploadItem(
                relative_path="PDF/security.pdf",
                content=_pdf_bytes("Security reviews are required quarterly."),
            ),
            FolderUploadItem(
                relative_path="Office/就業規則.docx",
                content=_docx_bytes("休暇制度", "有給休暇の申請手順です。"),
            ),
            FolderUploadItem(
                relative_path="Office/勤務表.xlsx",
                content=_xlsx_bytes(),
            ),
            FolderUploadItem(
                relative_path="Office/説明会.pptx",
                content=_pptx_bytes(),
            ),
            FolderUploadItem(
                relative_path="Office/broken.docx",
                content=b"not a valid docx file",
            ),
        ],
        uploader_user_id=7,
    )

    with session_factory() as session:
        documents = session.scalars(
            select(DocumentRecord).order_by(DocumentRecord.storage_path)
        ).all()
        chunks = session.scalars(select(DocumentChunkRecord)).all()
        permissions = session.scalars(select(DocumentPermissionRecord)).all()

    assert response.files_discovered == 6
    assert response.indexed == 5
    assert response.failed == 1
    assert response.results[-1].status == "failed"
    assert response.results[-1].reason == "extraction_failed"
    assert {document.file_type for document in documents} == {
        "markdown",
        "pdf",
        "docx",
        "xlsx",
        "pptx",
    }
    assert "Office/就業規則.docx" in {document.storage_path for document in documents}
    assert len(permissions) == 5
    assert vector_store.active_point_ids == {
        chunk.qdrant_point_id for chunk in chunks
    }


def test_folder_ingestion_skips_unsupported_hidden_binary_large_and_unsafe_paths(
    tmp_path,
) -> None:
    ingestion_service, *_ = _build_folder_services(tmp_path, max_upload_bytes=20)

    response = ingestion_service.ingest_folder_documents(
        folder_name="/Users/uploader/CompanyDocs",
        files=[
            FolderUploadItem(relative_path=".DS_Store", content=b"metadata"),
            FolderUploadItem(relative_path=".git/config", content=b"[core]\n"),
            FolderUploadItem(
                relative_path="node_modules/pkg/readme.md",
                content=b"# Package",
            ),
            FolderUploadItem(relative_path="notes.txt", content=b"plain text"),
            FolderUploadItem(relative_path="HR/binary.md", content=b"\x00\x01\x02"),
            FolderUploadItem(relative_path="HR/big.md", content=b"#" * 21),
            FolderUploadItem(relative_path="../secret.md", content=b"# Secret"),
        ],
        uploader_user_id=7,
    )

    assert response.folder_name == "CompanyDocs"
    assert response.indexed == 0
    assert response.skipped == 7
    assert response.failed == 0
    assert response.skip_reasons == {
        "hidden_or_system_file": 1,
        "excluded_directory": 2,
        "unsupported_extension": 1,
        "binary_file": 1,
        "too_large": 1,
        "unsafe_path": 1,
    }


def test_folder_ingestion_returns_duplicates_as_skipped_without_extra_vectors(
    tmp_path,
) -> None:
    (
        ingestion_service,
        _metadata_service,
        permission_service,
        vector_store,
        session_factory,
    ) = _build_folder_services(tmp_path)
    content = b"# Security Policy\n\nUse hardware keys for admin access."

    first_response = ingestion_service.ingest_folder_documents(
        folder_name="CompanyDocs",
        files=[FolderUploadItem(relative_path="IT/security.md", content=content)],
        uploader_user_id=7,
    )
    second_response = ingestion_service.ingest_folder_documents(
        folder_name="CompanyDocs",
        files=[
            FolderUploadItem(relative_path="IT/security.md", content=content),
            FolderUploadItem(relative_path="Copies/security.md", content=content),
        ],
        uploader_user_id=8,
    )

    with session_factory() as session:
        documents = session.scalars(select(DocumentRecord)).all()
        chunks = session.scalars(select(DocumentChunkRecord)).all()

    assert first_response.indexed == 1
    assert second_response.indexed == 0
    assert second_response.skipped == 2
    assert {result.reason for result in second_response.results} == {
        "already_indexed"
    }
    assert len(documents) == 1
    assert len(chunks) == 1
    assert len(vector_store.stored_point_ids) == 1
    assert vector_store.active_point_ids == set(vector_store.stored_point_ids)
    assert permission_service.list_accessible_document_ids(8) == [
        first_response.results[0].document_id
    ]


def test_folder_ingestion_continues_after_file_failure_and_cleans_partial_state(
    tmp_path,
) -> None:
    engine = create_engine("sqlite:///:memory:")
    Base.metadata.create_all(bind=engine)
    session_factory = sessionmaker(bind=engine, expire_on_commit=False)
    metadata_service = FailingForPathMetadataService(
        session_factory=session_factory,
        init_database=lambda: Base.metadata.create_all(bind=engine),
        fail_source_path="HR/broken.md",
    )
    permission_service = PermissionService(
        session_factory=session_factory,
        init_database=lambda: Base.metadata.create_all(bind=engine),
    )
    with session_factory() as session:
        session.add(
            UserRecord(
                id=7,
                email="uploader@example.com",
                password_hash="$argon2id$hash",
            )
        )
        session.commit()
    vector_store = TrackingVectorStore()
    ingestion_service = IngestionService(
        documents_dir=tmp_path,
        chunk_config=ChunkingConfig(chunk_size=200, chunk_overlap=20),
        embedding_service=FakeEmbeddingService(),
        vector_store=vector_store,
        metadata_service=metadata_service,
        permission_service=permission_service,
    )

    response = ingestion_service.ingest_folder_documents(
        folder_name="CompanyDocs",
        files=[
            FolderUploadItem(
                relative_path="HR/broken.md",
                content=b"# Broken\n\nThis fails after vectors are stored.",
            ),
            FolderUploadItem(
                relative_path="IT/security.md",
                content=b"# Security Policy\n\nUse hardware keys for admin access.",
            ),
        ],
        uploader_user_id=7,
    )

    with session_factory() as session:
        documents = session.scalars(select(DocumentRecord)).all()
        chunks = session.scalars(select(DocumentChunkRecord)).all()
        permissions = session.scalars(select(DocumentPermissionRecord)).all()

    assert response.indexed == 1
    assert response.failed == 1
    assert [result.status for result in response.results] == ["failed", "indexed"]
    assert response.results[0].reason == "metadata_error"
    assert [document.storage_path for document in documents] == ["IT/security.md"]
    assert len(permissions) == 1
    assert not (tmp_path / "HR" / "broken.md").exists()
    assert vector_store.deleted_point_ids
    assert vector_store.active_point_ids == {
        chunk.qdrant_point_id for chunk in chunks
    }


def test_folder_ingestion_reports_corrupt_pdf_without_blocking_good_files(
    tmp_path,
) -> None:
    ingestion_service, *_ = _build_folder_services(tmp_path)

    response = ingestion_service.ingest_folder_documents(
        folder_name="MixedDocs",
        files=[
            FolderUploadItem(
                relative_path="good.md",
                content=b"# Good\n\nThis markdown document is indexed.",
            ),
            FolderUploadItem(
                relative_path="broken.pdf",
                content=b"this is not a pdf",
            ),
        ],
        uploader_user_id=7,
    )

    assert response.indexed == 1
    assert response.failed == 1
    assert response.results[0].status == "indexed"
    assert response.results[1].status == "failed"
    assert response.results[1].reason == "unsupported_or_corrupt_pdf"


def _docx_bytes(heading: str, body: str) -> bytes:
    buffer = BytesIO()
    document = DocxDocument()
    document.add_heading(heading, level=1)
    document.add_paragraph(body)
    document.save(buffer)
    return buffer.getvalue()


def _xlsx_bytes() -> bytes:
    buffer = BytesIO()
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "勤怠"
    sheet.append(["社員", "状態"])
    sheet.append(["山田太郎", "在宅勤務"])
    workbook.create_sheet("空シート")
    workbook.save(buffer)
    return buffer.getvalue()


def _pptx_bytes() -> bytes:
    buffer = BytesIO()
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text = "全社会議"
    slide.placeholders[1].text = "検索品質を確認します。"
    presentation.save(buffer)
    return buffer.getvalue()


def _pdf_bytes(text: str) -> bytes:
    buffer = BytesIO()
    document = canvas.Canvas(buffer, pagesize=letter)
    document.setFont("Helvetica", 12)
    document.drawString(72, 720, text)
    document.save()
    return buffer.getvalue()
