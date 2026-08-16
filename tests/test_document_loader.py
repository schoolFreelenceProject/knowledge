from pathlib import Path

import pytest
from docx import Document as DocxDocument
from openpyxl import Workbook
from PIL import Image, ImageDraw
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.cidfonts import UnicodeCIDFont

from app.services.document_loader import (
    CorruptPdfError,
    DocumentLoaderError,
    PdfExtractionConfig,
    PdfExtractionError,
    ScannedPdfRequiresOcrError,
    load_document,
)


def test_load_text_based_pdf_extracts_text_and_page_metadata(tmp_path) -> None:
    path = tmp_path / "policy.pdf"
    _write_text_pdf(path, "Remote work policy applies to all employees.")

    documents = load_document(
        path,
        documents_dir=tmp_path,
        pdf_config=PdfExtractionConfig(min_text_chars=10, ocr_enabled=False),
    )

    assert len(documents) == 1
    assert "Remote work policy" in documents[0].text
    assert documents[0].metadata.filename == "policy.pdf"
    assert documents[0].metadata.source_path == "policy.pdf"
    assert documents[0].metadata.page_number == 1


def test_load_japanese_text_pdf_extracts_selectable_japanese_text(tmp_path) -> None:
    path = tmp_path / "japanese.pdf"
    japanese_text = "日本語能力試験N1の読解問題です。語彙と文法を確認します。"
    _write_japanese_pdf(path, japanese_text)

    documents = load_document(
        path,
        documents_dir=tmp_path,
        pdf_config=PdfExtractionConfig(min_text_chars=10, ocr_enabled=False),
    )

    assert len(documents) == 1
    assert "日本語能力試験" in documents[0].text
    assert documents[0].metadata.page_number == 1


def test_load_japanese_docx_extracts_sections_lists_tables_and_metadata(
    tmp_path,
) -> None:
    path = tmp_path / "就業規則.docx"
    document = DocxDocument()
    document.add_heading("休暇制度", level=1)
    document.add_paragraph("ＶＰＮ　接続中でも休暇申請を確認できます。")
    document.add_paragraph("有給休暇の残日数を確認する", style="List Bullet")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "項目"
    table.cell(0, 1).text = "説明"
    table.cell(1, 0).text = "経費精算"
    table.cell(1, 1).text = "領収書を添付します"
    document.save(path)

    documents = load_document(path, documents_dir=tmp_path)

    assert len(documents) == 1
    assert "VPN 接続中" in documents[0].text
    assert "- 有給休暇" in documents[0].text
    assert "経費精算" in documents[0].text
    assert documents[0].metadata.file_type == "docx"
    assert documents[0].metadata.source_path == "就業規則.docx"
    assert documents[0].metadata.section_heading == "休暇制度"
    assert documents[0].metadata.heading_path == "休暇制度"
    assert "table" in (documents[0].metadata.block_kind or "")


def test_load_japanese_xlsx_extracts_sheets_cells_and_ranges(tmp_path) -> None:
    path = tmp_path / "勤務表.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "勤怠"
    sheet.append(["社員", "状態", "備考"])
    sheet.append(["山田太郎", "在宅勤務", "VPN 接続"])
    glossary = workbook.create_sheet("用語")
    glossary.append(["単語", "説明"])
    glossary.append(["経費精算", "領収書を添付する"])
    workbook.create_sheet("空シート")
    workbook.save(path)

    documents = load_document(path, documents_dir=tmp_path)

    assert [document.metadata.sheet_name for document in documents] == ["勤怠", "用語"]
    assert documents[0].metadata.workbook == "勤務表.xlsx"
    assert documents[0].metadata.file_type == "xlsx"
    assert documents[0].metadata.cell_range == "A1:C2"
    assert documents[0].metadata.row_start == 1
    assert documents[0].metadata.row_end == 2
    assert "山田太郎" in documents[0].text
    assert "経費精算" in documents[1].text


def test_load_japanese_pptx_extracts_slide_text_notes_and_metadata(tmp_path) -> None:
    path = tmp_path / "説明会.pptx"
    presentation = Presentation()
    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = "全社会議"
    title_slide.placeholders[1].text = "品質保証ロードマップ"
    _set_slide_notes(title_slide, "発表者ノート: 日本語OCR結果を確認")

    second_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    second_slide.shapes.title.text = "検索改善"
    second_slide.placeholders[1].text = "Kanji / Hiragana / カタカナ / English"
    presentation.save(path)

    documents = load_document(path, documents_dir=tmp_path)

    assert len(documents) == 2
    assert documents[0].metadata.file_type == "pptx"
    assert documents[0].metadata.slide_number == 1
    assert documents[0].metadata.slide_title == "全社会議"
    assert "品質保証ロードマップ" in documents[0].text
    assert "日本語OCR結果" in documents[0].text
    assert documents[1].metadata.slide_number == 2
    assert documents[1].metadata.slide_title == "検索改善"


def test_load_scanned_pdf_requires_ocr_when_ocr_is_disabled(tmp_path) -> None:
    path = tmp_path / "scanned.pdf"
    _write_image_only_pdf(path)

    with pytest.raises(ScannedPdfRequiresOcrError):
        load_document(
            path,
            documents_dir=tmp_path,
            pdf_config=PdfExtractionConfig(min_text_chars=10, ocr_enabled=False),
        )


def test_load_scanned_pdf_uses_ocr_only_for_insufficient_pages(
    tmp_path,
    monkeypatch,
) -> None:
    path = tmp_path / "scanned.pdf"
    _write_image_only_pdf(path)
    calls: list[int] = []

    def fake_ocr_pdf_page(path: Path, page_number: int, config: PdfExtractionConfig):
        calls.append(page_number)
        return "日本語 OCR テキスト 読解 文法 語彙"

    monkeypatch.setattr(
        "app.services.document_loader._ocr_pdf_page",
        fake_ocr_pdf_page,
    )

    documents = load_document(
        path,
        documents_dir=tmp_path,
        pdf_config=PdfExtractionConfig(min_text_chars=10, ocr_enabled=True),
    )

    assert calls == [1]
    assert len(documents) == 1
    assert "日本語 OCR" in documents[0].text


def test_load_scanned_pdf_reports_extraction_failed_when_ocr_has_no_text(
    tmp_path,
    monkeypatch,
) -> None:
    path = tmp_path / "qr-only.pdf"
    _write_image_only_pdf(path)

    monkeypatch.setattr(
        "app.services.document_loader._ocr_pdf_page",
        lambda path, page_number, config: "",
    )

    with pytest.raises(PdfExtractionError, match="OCR produced no usable text"):
        load_document(
            path,
            documents_dir=tmp_path,
            pdf_config=PdfExtractionConfig(min_text_chars=10, ocr_enabled=True),
        )


def test_load_corrupt_pdf_reports_corrupt_pdf(tmp_path) -> None:
    path = tmp_path / "broken.pdf"
    path.write_bytes(b"this is not a pdf")

    with pytest.raises(CorruptPdfError):
        load_document(
            path,
            documents_dir=tmp_path,
            pdf_config=PdfExtractionConfig(min_text_chars=10, ocr_enabled=False),
        )


@pytest.mark.parametrize("extension", [".docx", ".xlsx", ".pptx"])
def test_load_corrupt_office_file_reports_loader_error(tmp_path, extension) -> None:
    path = tmp_path / f"broken{extension}"
    path.write_bytes(b"this is not a valid office document")

    with pytest.raises(DocumentLoaderError):
        load_document(path, documents_dir=tmp_path)


def _write_text_pdf(path: Path, text: str) -> None:
    document = canvas.Canvas(str(path), pagesize=letter)
    document.setFont("Helvetica", 12)
    document.drawString(72, 720, text)
    document.save()


def _write_japanese_pdf(path: Path, text: str) -> None:
    font_name = "HeiseiMin-W3"
    pdfmetrics.registerFont(UnicodeCIDFont(font_name))
    document = canvas.Canvas(str(path), pagesize=letter)
    document.setFont(font_name, 12)
    document.drawString(72, 720, text)
    document.save()


def _write_image_only_pdf(path: Path) -> None:
    image = Image.new("RGB", (640, 180), "white")
    draw = ImageDraw.Draw(image)
    draw.text((32, 72), "Scanned PDF page", fill="black")

    document = canvas.Canvas(str(path), pagesize=letter)
    document.drawInlineImage(image, 72, 580, width=320, height=90)
    document.save()


def _set_slide_notes(slide, text: str) -> None:
    try:
        slide.notes_slide.notes_text_frame.text = text
    except Exception:
        pass
