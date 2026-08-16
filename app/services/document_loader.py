import re
import shutil
import subprocess
import tempfile
from dataclasses import dataclass
from datetime import date, datetime, time
from decimal import Decimal
from pathlib import Path
from typing import Any

from app.schemas.documents import DocumentFileType, DocumentMetadata, ExtractedDocument
from app.services.text_normalization import normalize_text


PDF_EXTENSION = ".pdf"
MARKDOWN_EXTENSIONS = {".md", ".markdown"}
DOCX_EXTENSION = ".docx"
XLSX_EXTENSION = ".xlsx"
PPTX_EXTENSION = ".pptx"
OFFICE_EXTENSIONS = {DOCX_EXTENSION, XLSX_EXTENSION, PPTX_EXTENSION}
SUPPORTED_EXTENSIONS = {PDF_EXTENSION, *MARKDOWN_EXTENSIONS, *OFFICE_EXTENSIONS}
DEFAULT_PDF_MIN_TEXT_CHARS = 20
DEFAULT_PDF_OCR_DPI = 200
DEFAULT_PDF_OCR_LANGUAGES = "jpn+eng"
DEFAULT_PDF_OCR_TIMEOUT_SECONDS = 120
DEFAULT_PDF_OCR_MAX_PAGES = 100
DEFAULT_PDF_TEXT_EXTRACTION_TIMEOUT_SECONDS = 30
XLSX_ROWS_PER_BLOCK = 100
XLSX_MAX_CHARS_PER_BLOCK = 12_000


@dataclass(frozen=True)
class PdfExtractionConfig:
    min_text_chars: int = DEFAULT_PDF_MIN_TEXT_CHARS
    ocr_enabled: bool = False
    ocr_languages: str = DEFAULT_PDF_OCR_LANGUAGES
    ocr_dpi: int = DEFAULT_PDF_OCR_DPI
    ocr_timeout_seconds: int = DEFAULT_PDF_OCR_TIMEOUT_SECONDS
    ocr_max_pages: int = DEFAULT_PDF_OCR_MAX_PAGES
    text_extraction_timeout_seconds: int = (
        DEFAULT_PDF_TEXT_EXTRACTION_TIMEOUT_SECONDS
    )


class DocumentLoaderError(RuntimeError):
    """Raised when a document cannot be loaded or parsed."""


class UnsupportedDocumentTypeError(DocumentLoaderError):
    """Raised when a file extension is not supported by the MVP parser."""


class PdfExtractionError(DocumentLoaderError):
    """Raised when PDF text cannot be extracted into usable content."""


class ScannedPdfRequiresOcrError(PdfExtractionError):
    """Raised when a PDF appears image-only and OCR is disabled or unavailable."""


class EncryptedPdfError(PdfExtractionError):
    """Raised when a PDF is encrypted and cannot be read without a password."""


class CorruptPdfError(PdfExtractionError):
    """Raised when a PDF cannot be parsed as a supported PDF file."""


def discover_document_paths(documents_dir: str | Path) -> list[Path]:
    root = Path(documents_dir)
    if not root.exists():
        raise FileNotFoundError(f"Documents directory does not exist: {root}")
    if not root.is_dir():
        raise NotADirectoryError(f"Documents path is not a directory: {root}")

    return sorted(
        path
        for path in root.rglob("*")
        if path.is_file() and path.suffix.lower() in SUPPORTED_EXTENSIONS
    )


def load_documents(documents_dir: str | Path) -> list[ExtractedDocument]:
    root = Path(documents_dir)
    extracted_documents: list[ExtractedDocument] = []

    for document_path in discover_document_paths(root):
        extracted_documents.extend(load_document(document_path, documents_dir=root))

    return extracted_documents


def load_document(
    document_path: str | Path,
    documents_dir: str | Path | None = None,
    pdf_config: PdfExtractionConfig | None = None,
) -> list[ExtractedDocument]:
    path = Path(document_path)
    if not path.exists():
        raise FileNotFoundError(f"Document does not exist: {path}")
    if not path.is_file():
        raise DocumentLoaderError(f"Document path is not a file: {path}")

    base_dir = Path(documents_dir) if documents_dir is not None else path.parent
    extension = path.suffix.lower()

    if extension == PDF_EXTENSION:
        return _load_pdf(path, base_dir=base_dir, config=_resolve_pdf_config(pdf_config))
    if extension in MARKDOWN_EXTENSIONS:
        return [_load_markdown(path, base_dir=base_dir)]
    if extension == DOCX_EXTENSION:
        return _load_docx(path, base_dir=base_dir)
    if extension == XLSX_EXTENSION:
        return _load_xlsx(path, base_dir=base_dir)
    if extension == PPTX_EXTENSION:
        return _load_pptx(path, base_dir=base_dir)

    raise UnsupportedDocumentTypeError(
        f"Unsupported document type '{extension}'. "
        f"Supported types: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
    )


def _load_pdf(
    path: Path,
    base_dir: Path,
    config: PdfExtractionConfig,
) -> list[ExtractedDocument]:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise DocumentLoaderError(
            "PDF parsing requires pypdf. Install dependencies with "
            "`pip install -r requirements.txt`."
        ) from exc

    try:
        reader = PdfReader(str(path))
    except Exception as exc:  # pypdf exposes multiple parse-time exceptions.
        raise CorruptPdfError(f"Failed to read PDF '{path}'.") from exc

    if reader.is_encrypted and not _decrypt_pdf_without_password(reader):
        raise EncryptedPdfError(f"PDF '{path}' is encrypted.")

    try:
        page_count = len(reader.pages)
    except Exception as exc:
        raise CorruptPdfError(f"Failed to inspect PDF pages for '{path}'.") from exc

    page_texts: list[str] = []
    for page in reader.pages:
        try:
            page_texts.append(_normalize_text(page.extract_text() or ""))
        except Exception:
            page_texts.append("")

    _fill_insufficient_pages_with_poppler_text(path=path, page_texts=page_texts, config=config)
    insufficient_page_indexes = _insufficient_page_indexes(page_texts, config)
    if insufficient_page_indexes:
        if config.ocr_enabled:
            if len(insufficient_page_indexes) > config.ocr_max_pages:
                raise ScannedPdfRequiresOcrError(
                    "PDF requires OCR for too many pages. "
                    f"Maximum OCR page count is {config.ocr_max_pages}."
                )
            _fill_insufficient_pages_with_ocr(
                path=path,
                page_texts=page_texts,
                page_indexes=insufficient_page_indexes,
                config=config,
            )
        elif _has_no_usable_text(page_texts, config):
            raise ScannedPdfRequiresOcrError(
                "PDF appears to be scanned or image-only and requires OCR."
            )

    extracted_pages: list[ExtractedDocument] = []
    for page_index, text in enumerate(page_texts, start=1):
        if _meaningful_text_length(text) < config.min_text_chars:
            continue
        metadata = _build_metadata(
            path=path,
            base_dir=base_dir,
            file_type="pdf",
            page_number=page_index,
        )
        extracted_pages.append(ExtractedDocument(text=text, metadata=metadata))

    if not extracted_pages:
        if _pdf_has_images(path):
            if config.ocr_enabled:
                raise PdfExtractionError("PDF OCR produced no usable text.")
            raise ScannedPdfRequiresOcrError(
                "PDF appears to be scanned or image-only and requires OCR."
            )
        raise PdfExtractionError("PDF text extraction produced no usable text.")

    return extracted_pages


def _load_markdown(path: Path, base_dir: Path) -> ExtractedDocument:
    try:
        text = path.read_text(encoding="utf-8-sig")
    except UnicodeDecodeError as exc:
        raise DocumentLoaderError(
            f"Failed to read Markdown file '{path}' as UTF-8 text."
        ) from exc

    metadata = _build_metadata(
        path=path,
        base_dir=base_dir,
        file_type="markdown",
        page_number=None,
    )
    return ExtractedDocument(text=_normalize_text(text), metadata=metadata)


def _load_docx(path: Path, base_dir: Path) -> list[ExtractedDocument]:
    try:
        from docx import Document as DocxDocument
        from docx.table import Table
        from docx.text.paragraph import Paragraph
    except ImportError as exc:
        raise DocumentLoaderError(
            "DOCX parsing requires python-docx. Install dependencies with "
            "`pip install -r requirements.txt`."
        ) from exc

    try:
        document = DocxDocument(str(path))
    except Exception as exc:
        raise DocumentLoaderError(f"Failed to read DOCX document '{path}'.") from exc

    extracted_sections: list[ExtractedDocument] = []
    headings: list[str] = []
    current_lines: list[str] = []
    current_heading: str | None = None
    current_block_kinds: set[str] = set()
    table_index = 0

    def flush_section() -> None:
        nonlocal current_lines, current_heading, current_block_kinds
        text = _normalize_text("\n".join(current_lines))
        if text:
            extracted_sections.append(
                ExtractedDocument(
                    text=text,
                    metadata=_build_metadata(
                        path=path,
                        base_dir=base_dir,
                        file_type="docx",
                        page_number=None,
                        section_heading=current_heading,
                        heading_path=" > ".join(headings) if headings else None,
                        block_kind=(
                            "+".join(sorted(current_block_kinds))
                            if current_block_kinds
                            else "section"
                        ),
                    ),
                )
            )
        current_lines = []
        current_heading = headings[-1] if headings else None
        current_block_kinds = set()

    for block in _iter_docx_blocks(
        document=document,
        paragraph_type=Paragraph,
        table_type=Table,
    ):
        if isinstance(block, Paragraph):
            paragraph_text = _normalize_text(block.text)
            if not paragraph_text:
                continue

            heading_level = _docx_heading_level(block)
            if heading_level is not None:
                flush_section()
                headings = _update_heading_path(
                    headings=headings,
                    level=heading_level,
                    heading=paragraph_text,
                )
                current_heading = paragraph_text
                current_lines.append(f"{'#' * min(heading_level, 6)} {paragraph_text}")
                current_block_kinds.add("heading")
                continue

            if _is_docx_list_paragraph(block):
                current_lines.append(f"- {paragraph_text}")
                current_block_kinds.add("list")
            else:
                current_lines.append(paragraph_text)
                current_block_kinds.add("paragraph")
            continue

        table_index += 1
        table_text = _docx_table_text(block, table_index=table_index)
        if table_text:
            current_lines.append(table_text)
            current_block_kinds.add("table")

    flush_section()
    if not extracted_sections:
        raise DocumentLoaderError("DOCX text extraction produced no usable text.")

    return extracted_sections


def _load_xlsx(path: Path, base_dir: Path) -> list[ExtractedDocument]:
    try:
        from openpyxl import load_workbook
        from openpyxl.utils import get_column_letter
    except ImportError as exc:
        raise DocumentLoaderError(
            "XLSX parsing requires openpyxl. Install dependencies with "
            "`pip install -r requirements.txt`."
        ) from exc

    try:
        workbook = load_workbook(
            filename=path,
            read_only=True,
            data_only=True,
        )
    except Exception as exc:
        raise DocumentLoaderError(f"Failed to read XLSX workbook '{path}'.") from exc

    extracted_blocks: list[ExtractedDocument] = []
    try:
        for worksheet in workbook.worksheets:
            sheet_blocks = _extract_xlsx_sheet_blocks(
                path=path,
                base_dir=base_dir,
                workbook_name=path.name,
                sheet=worksheet,
                get_column_letter=get_column_letter,
            )
            extracted_blocks.extend(sheet_blocks)
    finally:
        workbook.close()

    if not extracted_blocks:
        raise DocumentLoaderError("XLSX text extraction produced no usable text.")

    return extracted_blocks


def _load_pptx(path: Path, base_dir: Path) -> list[ExtractedDocument]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise DocumentLoaderError(
            "PPTX parsing requires python-pptx. Install dependencies with "
            "`pip install -r requirements.txt`."
        ) from exc

    try:
        presentation = Presentation(str(path))
    except Exception as exc:
        raise DocumentLoaderError(f"Failed to read PPTX deck '{path}'.") from exc

    extracted_slides: list[ExtractedDocument] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_title = _pptx_slide_title(slide)
        slide_text_parts = [
            text
            for text in (_pptx_shape_text(shape) for shape in slide.shapes)
            if text
        ]
        notes_text = _pptx_notes_text(slide)
        if not slide_title and not slide_text_parts and not notes_text:
            continue

        lines = [f"Slide {slide_number}"]
        if slide_title:
            lines[0] = f"Slide {slide_number}: {slide_title}"
        lines.extend(slide_text_parts)
        if notes_text:
            lines.extend(["Speaker notes:", notes_text])

        text = _normalize_text("\n".join(lines))
        if _meaningful_text_length(text) == 0:
            continue

        extracted_slides.append(
            ExtractedDocument(
                text=text,
                metadata=_build_metadata(
                    path=path,
                    base_dir=base_dir,
                    file_type="pptx",
                    page_number=None,
                    slide_number=slide_number,
                    slide_title=slide_title,
                    block_kind="slide",
                ),
            )
        )

    if not extracted_slides:
        raise DocumentLoaderError("PPTX text extraction produced no usable text.")

    return extracted_slides


def _build_metadata(
    path: Path,
    base_dir: Path,
    file_type: DocumentFileType,
    page_number: int | None,
    **extra_metadata: Any,
) -> DocumentMetadata:
    source_path = _relative_source_path(path=path, base_dir=base_dir)
    return DocumentMetadata(
        filename=source_path if "/" in source_path else path.name,
        source_path=source_path,
        file_type=file_type,
        page_number=page_number,
        **extra_metadata,
    )


def _iter_docx_blocks(document, paragraph_type, table_type):
    for child in document.element.body.iterchildren():
        if child.tag.endswith("}p"):
            yield paragraph_type(child, document)
        elif child.tag.endswith("}tbl"):
            yield table_type(child, document)


def _docx_heading_level(paragraph) -> int | None:
    style_name = str(getattr(getattr(paragraph, "style", None), "name", "") or "")
    normalized_style = style_name.casefold()
    is_heading = (
        "heading" in normalized_style
        or "見出し" in style_name
        or normalized_style == "title"
    )
    if not is_heading:
        return None

    level_match = re.search(r"\d+", style_name)
    if level_match is None:
        return 1

    return max(1, min(int(level_match.group(0)), 6))


def _update_heading_path(
    headings: list[str],
    level: int,
    heading: str,
) -> list[str]:
    safe_level = max(1, level)
    updated = headings[: safe_level - 1]
    updated.append(heading)
    return updated


def _is_docx_list_paragraph(paragraph) -> bool:
    style_name = str(getattr(getattr(paragraph, "style", None), "name", "") or "")
    normalized_style = style_name.casefold()
    if any(
        marker in normalized_style
        for marker in ("list", "bullet", "number")
    ) or "箇条書き" in style_name:
        return True

    try:
        paragraph_properties = paragraph._p.pPr
    except AttributeError:
        return False

    return (
        paragraph_properties is not None
        and getattr(paragraph_properties, "numPr", None) is not None
    )


def _docx_table_text(table, table_index: int) -> str:
    rows: list[str] = []
    for row in table.rows:
        cells = [
            _normalize_text(cell.text).replace("\n", " / ")
            for cell in row.cells
        ]
        if any(cells):
            rows.append(" | ".join(cells).rstrip())

    if not rows:
        return ""

    return _normalize_text("\n".join([f"Table {table_index}", *rows]))


def _extract_xlsx_sheet_blocks(
    path: Path,
    base_dir: Path,
    workbook_name: str,
    sheet,
    get_column_letter,
) -> list[ExtractedDocument]:
    extracted_blocks: list[ExtractedDocument] = []
    sheet_name = _normalize_text(str(sheet.title))
    current_rows: list[str] = []
    current_char_count = 0
    row_start: int | None = None
    row_end: int | None = None
    min_col: int | None = None
    max_col: int | None = None

    def flush_block() -> None:
        nonlocal current_rows, current_char_count, row_start, row_end, min_col, max_col
        if (
            not current_rows
            or row_start is None
            or row_end is None
            or min_col is None
            or max_col is None
        ):
            current_rows = []
            current_char_count = 0
            row_start = None
            row_end = None
            min_col = None
            max_col = None
            return

        cell_range = (
            f"{get_column_letter(min_col)}{row_start}:"
            f"{get_column_letter(max_col)}{row_end}"
        )
        text = _normalize_text(
            "\n".join(
                [
                    f"Workbook: {workbook_name}",
                    f"Sheet: {sheet_name}",
                    f"Rows: {row_start}-{row_end}",
                    f"Cell range: {cell_range}",
                    *current_rows,
                ]
            )
        )
        if text:
            extracted_blocks.append(
                ExtractedDocument(
                    text=text,
                    metadata=_build_metadata(
                        path=path,
                        base_dir=base_dir,
                        file_type="xlsx",
                        page_number=None,
                        workbook=workbook_name,
                        sheet_name=sheet_name,
                        cell_range=cell_range,
                        row_start=row_start,
                        row_end=row_end,
                        block_kind="sheet_rows",
                    ),
                )
            )

        current_rows = []
        current_char_count = 0
        row_start = None
        row_end = None
        min_col = None
        max_col = None

    for row_number, row in enumerate(sheet.iter_rows(), start=1):
        row_cells: list[str] = []
        row_min_col: int | None = None
        row_max_col: int | None = None
        for column_number, cell in enumerate(row, start=1):
            cell_text = _format_xlsx_cell_value(getattr(cell, "value", None))
            if not cell_text:
                continue

            column_name = get_column_letter(column_number)
            row_cells.append(f"{column_name}={cell_text}")
            row_min_col = column_number if row_min_col is None else min(
                row_min_col,
                column_number,
            )
            row_max_col = column_number if row_max_col is None else max(
                row_max_col,
                column_number,
            )

        if not row_cells or row_min_col is None or row_max_col is None:
            continue

        row_text = f"Row {row_number}: {' | '.join(row_cells)}"
        row_would_overflow = (
            current_rows
            and current_char_count + len(row_text) > XLSX_MAX_CHARS_PER_BLOCK
        )
        if len(current_rows) >= XLSX_ROWS_PER_BLOCK or row_would_overflow:
            flush_block()

        if row_start is None:
            row_start = row_number
        row_end = row_number
        min_col = row_min_col if min_col is None else min(min_col, row_min_col)
        max_col = row_max_col if max_col is None else max(max_col, row_max_col)
        current_rows.append(row_text)
        current_char_count += len(row_text)

    flush_block()
    return extracted_blocks


def _format_xlsx_cell_value(value: Any) -> str:
    if value is None:
        return ""

    if isinstance(value, bool):
        return "TRUE" if value else "FALSE"
    if isinstance(value, (datetime, date, time)):
        return value.isoformat()
    if isinstance(value, Decimal):
        return str(value.normalize())

    return _normalize_text(str(value)).replace("\n", " ")


def _pptx_slide_title(slide) -> str | None:
    title_shape = getattr(getattr(slide, "shapes", None), "title", None)
    if title_shape is None:
        return None

    return _normalize_text(getattr(title_shape, "text", "")) or None


def _pptx_shape_text(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""

    return _normalize_text(getattr(shape, "text", ""))


def _pptx_notes_text(slide) -> str | None:
    try:
        if not getattr(slide, "has_notes_slide", False):
            return None
        notes_slide = slide.notes_slide
        notes_frame = getattr(notes_slide, "notes_text_frame", None)
        if notes_frame is None:
            return None
    except Exception:
        return None

    return _normalize_text(getattr(notes_frame, "text", "")) or None


def _relative_source_path(path: Path, base_dir: Path) -> str:
    resolved_path = path.resolve()
    resolved_base_dir = base_dir.resolve()

    try:
        return resolved_path.relative_to(resolved_base_dir).as_posix()
    except ValueError:
        return resolved_path.as_posix()


def _normalize_text(text: str) -> str:
    return normalize_text(text)


def _resolve_pdf_config(
    config: PdfExtractionConfig | None,
) -> PdfExtractionConfig:
    return config or PdfExtractionConfig()


def _decrypt_pdf_without_password(reader) -> bool:
    try:
        return bool(reader.decrypt(""))
    except Exception:
        return False


def _fill_insufficient_pages_with_poppler_text(
    path: Path,
    page_texts: list[str],
    config: PdfExtractionConfig,
) -> None:
    for page_index in _insufficient_page_indexes(page_texts, config):
        poppler_text = _extract_pdf_page_text_with_poppler(
            path=path,
            page_number=page_index + 1,
            timeout_seconds=config.text_extraction_timeout_seconds,
        )
        if poppler_text is None:
            continue
        if _meaningful_text_length(poppler_text) > _meaningful_text_length(
            page_texts[page_index]
        ):
            page_texts[page_index] = poppler_text


def _fill_insufficient_pages_with_ocr(
    path: Path,
    page_texts: list[str],
    page_indexes: list[int],
    config: PdfExtractionConfig,
) -> None:
    for page_index in page_indexes:
        ocr_text = _ocr_pdf_page(
            path=path,
            page_number=page_index + 1,
            config=config,
        )
        if _meaningful_text_length(ocr_text) > _meaningful_text_length(
            page_texts[page_index]
        ):
            page_texts[page_index] = ocr_text


def _extract_pdf_page_text_with_poppler(
    path: Path,
    page_number: int,
    timeout_seconds: int,
) -> str | None:
    if shutil.which("pdftotext") is None:
        return None

    try:
        completed = subprocess.run(
            [
                "pdftotext",
                "-layout",
                "-f",
                str(page_number),
                "-l",
                str(page_number),
                str(path),
                "-",
            ],
            check=False,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
            encoding="utf-8",
            errors="replace",
            timeout=timeout_seconds,
        )
    except subprocess.TimeoutExpired:
        return None
    if completed.returncode != 0:
        return None

    return _normalize_text(completed.stdout)


def _ocr_pdf_page(
    path: Path,
    page_number: int,
    config: PdfExtractionConfig,
) -> str:
    if shutil.which("pdftoppm") is None or shutil.which("tesseract") is None:
        raise ScannedPdfRequiresOcrError(
            "PDF appears to be scanned or image-only and OCR tools are unavailable."
        )

    with tempfile.TemporaryDirectory(prefix="pdf-ocr-") as temp_dir:
        image_prefix = Path(temp_dir) / "page"
        try:
            render = subprocess.run(
                [
                    "pdftoppm",
                    "-f",
                    str(page_number),
                    "-l",
                    str(page_number),
                    "-r",
                    str(config.ocr_dpi),
                    "-png",
                    str(path),
                    str(image_prefix),
                ],
                check=False,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                text=True,
                encoding="utf-8",
                errors="replace",
                timeout=config.ocr_timeout_seconds,
            )
        except subprocess.TimeoutExpired as exc:
            raise PdfExtractionError("PDF page rendering timed out for OCR.") from exc
        if render.returncode != 0:
            raise PdfExtractionError("PDF page could not be rendered for OCR.")

        image_paths = sorted(Path(temp_dir).glob("page-*.png"))
        if not image_paths:
            raise PdfExtractionError("PDF page rendering produced no OCR image.")

        try:
            ocr = subprocess.run(
                [
                    "tesseract",
                    str(image_paths[0]),
                    "stdout",
                    "-l",
                    config.ocr_languages,
                ],
                check=False,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                text=True,
                encoding="utf-8",
                errors="replace",
                timeout=config.ocr_timeout_seconds,
            )
        except subprocess.TimeoutExpired as exc:
            raise PdfExtractionError("PDF OCR timed out.") from exc
        if ocr.returncode != 0:
            if (
                "Error opening data file" in ocr.stderr
                or "Failed loading language" in ocr.stderr
            ):
                raise ScannedPdfRequiresOcrError(
                    "PDF appears to be scanned or image-only and OCR language data is unavailable."
                )
            raise PdfExtractionError("PDF OCR failed.")

        return _normalize_text(ocr.stdout)


def _insufficient_page_indexes(
    page_texts: list[str],
    config: PdfExtractionConfig,
) -> list[int]:
    return [
        index
        for index, text in enumerate(page_texts)
        if _meaningful_text_length(text) < config.min_text_chars
    ]


def _has_no_usable_text(
    page_texts: list[str],
    config: PdfExtractionConfig,
) -> bool:
    return all(
        _meaningful_text_length(text) < config.min_text_chars
        for text in page_texts
    )


def _meaningful_text_length(text: str) -> int:
    return sum(1 for character in text if character.isalnum())


def _pdf_has_images(path: Path) -> bool:
    if shutil.which("pdfimages") is None:
        return False

    try:
        completed = subprocess.run(
            ["pdfimages", "-list", str(path)],
            check=False,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
            encoding="utf-8",
            errors="replace",
            timeout=DEFAULT_PDF_TEXT_EXTRACTION_TIMEOUT_SECONDS,
        )
    except subprocess.TimeoutExpired:
        return False
    if completed.returncode != 0:
        return False

    return len(completed.stdout.splitlines()) > 2
